from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt


OUT = Path("figures/arcc_overview_editable.pptx")


COLORS = {
    "blue": ("DCEBFF", "2D67B1"),
    "orange": ("FFF0D8", "D68024"),
    "purple": ("E9E1FA", "6C55A3"),
    "green": ("DFF1D7", "4D8B3E"),
    "pink": ("F8D9E8", "C7487C"),
    "white": ("FFFFFF", "333333"),
    "arcc": ("FFF0DD", "C75524"),
    "line": ("333333", "333333"),
}


def rgb(hex_value):
    return RGBColor.from_string(hex_value)


def set_arrow(line_shape, end=True):
    ln = line_shape._element.spPr.get_or_add_ln()
    tag = "a:tailEnd" if not end else "a:headEnd"
    existing = ln.find("{http://schemas.openxmlformats.org/drawingml/2006/main}" + tag.split(":")[1])
    if existing is not None:
        ln.remove(existing)
    arrow = OxmlElement(tag)
    arrow.set("type", "triangle")
    ln.append(arrow)


def add_box(slide, x, y, w, h, text, style="white", font_size=11, bold=False, radius=True):
    fill, stroke = COLORS[style]
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = rgb(fill)
    shp.line.color.rgb = rgb(stroke)
    shp.line.width = Pt(1.2)
    tf = shp.text_frame
    tf.clear()
    tf.margin_left = Inches(0.06)
    tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = "Arial"
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(30, 30, 30)
    return shp


def add_label(slide, x, y, w, h, text, size=9, color="333333", italic=False, bold=False):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = "Arial"
    run.font.size = Pt(size)
    run.font.italic = italic
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    return box


def add_line(slide, x1, y1, x2, y2, dashed=False, width=1.1, arrow=True):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1),
        Inches(y1),
        Inches(x2),
        Inches(y2),
    )
    line.line.color.rgb = rgb("333333")
    line.line.width = Pt(width)
    if dashed:
        line.line.dash_style = 4
    if arrow:
        set_arrow(line, end=True)
    return line


def add_plus(slide, x, y):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(0.28), Inches(0.28))
    shp.fill.solid()
    shp.fill.fore_color.rgb = rgb("FFFFFF")
    shp.line.color.rgb = rgb("333333")
    shp.line.width = Pt(1)
    tf = shp.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "+"
    r.font.name = "Arial"
    r.font.size = Pt(16)
    r.font.bold = True
    return shp


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_label(slide, 3.55, 0.12, 6.2, 0.35, "Overview of the Proposed Framework", 20, "123A70", bold=True)

    # Main input and synthetic training path.
    add_box(slide, 0.25, 2.25, 1.25, 0.62, "Input normal\nmedical image", "white", 10, True)
    add_box(slide, 2.10, 2.12, 1.55, 0.95, "Frozen\nBiomedCLIP\nImage Encoder\n(shared)", "blue", 10, True)
    add_line(slide, 1.50, 2.56, 2.10, 2.56)
    add_label(slide, 1.42, 2.30, 0.75, 0.22, "main pass", 8, italic=True)

    add_box(slide, 1.75, 0.72, 1.85, 0.55, "Synthetic anomaly\ngeneration", "orange", 9, True)
    add_box(slide, 4.05, 0.72, 1.55, 0.55, "Synthetic\nimage + mask", "orange", 9, True)
    add_line(slide, 0.88, 2.25, 0.88, 0.98, dashed=True)
    add_line(slide, 0.88, 0.98, 1.75, 0.98, dashed=True)
    add_label(slide, 0.96, 0.58, 1.1, 0.22, "training only", 8, italic=True)
    add_line(slide, 3.60, 1.00, 4.05, 1.00)
    add_line(slide, 4.82, 1.27, 3.00, 2.12, dashed=True)
    add_label(slide, 3.36, 1.47, 1.2, 0.22, "synthetic image", 8, italic=True)

    # Image features.
    add_box(slide, 4.05, 1.78, 1.25, 0.45, "Global feature", "purple", 10)
    add_box(slide, 4.05, 2.83, 1.25, 0.45, "Patch features", "purple", 10)
    add_line(slide, 3.65, 2.38, 4.05, 2.00)
    add_line(slide, 3.65, 2.70, 4.05, 3.05)

    # Text branch.
    add_box(slide, 0.25, 5.58, 1.45, 0.58, "Normal / Abnormal\nprompts", "white", 10, True)
    add_box(slide, 2.08, 5.55, 1.45, 0.62, "Frozen BiomedCLIP\nText Encoder", "blue", 9, True)
    add_box(slide, 3.90, 5.55, 1.35, 0.62, "Learnable text\nprototypes", "orange", 9, True)
    add_line(slide, 1.70, 5.87, 2.08, 5.87)
    add_line(slide, 3.53, 5.87, 3.90, 5.87)

    # Global and patch image-level scoring.
    add_box(slide, 6.30, 1.68, 1.35, 0.58, "Global anomaly\nscoring", "orange", 10)
    add_box(slide, 8.20, 1.78, 1.08, 0.42, "Global score", "green", 10)
    add_line(slide, 5.30, 2.00, 6.30, 1.98)
    add_line(slide, 7.65, 1.98, 8.20, 1.98)

    add_box(slide, 5.45, 2.78, 1.25, 0.55, "Mamba /\nLocal Adapter", "orange", 9, True)
    add_box(slide, 7.05, 2.78, 1.42, 0.55, "Context-enhanced\npatch features", "purple", 9)
    add_box(slide, 8.90, 2.78, 1.38, 0.55, "Patch-level\nanomaly scoring", "orange", 9)
    add_box(slide, 10.65, 2.85, 0.9, 0.42, "Patch score", "green", 9)
    add_line(slide, 5.30, 3.05, 5.45, 3.05)
    add_line(slide, 6.70, 3.05, 7.05, 3.05)
    add_line(slide, 8.47, 3.05, 8.90, 3.05)
    add_line(slide, 10.28, 3.05, 10.65, 3.05)

    # Prototype guidance, as a simple bus.
    add_line(slide, 5.25, 5.55, 6.95, 2.26, arrow=False, width=0.9)
    add_line(slide, 6.95, 2.26, 6.95, 2.26)
    add_line(slide, 5.25, 5.55, 9.59, 2.78, arrow=False, width=0.9)
    add_label(slide, 5.75, 4.65, 1.2, 0.24, "text guidance", 8, italic=True)

    add_plus(slide, 11.90, 2.32)
    add_box(slide, 12.25, 2.18, 0.92, 0.58, "Image-level\nanomaly score", "green", 8.5)
    add_line(slide, 9.28, 1.98, 12.04, 2.46)
    add_line(slide, 11.55, 3.06, 12.04, 2.46)
    add_line(slide, 12.18, 2.46, 12.25, 2.46)

    # Pixel-level localization and ARCC.
    add_box(slide, 5.45, 4.28, 1.25, 0.55, "CNN localization\ndecoder", "orange", 9)
    add_box(slide, 7.05, 4.25, 1.20, 0.62, "Preliminary\nresponse\nA_local", "green", 8.5)
    add_line(slide, 4.68, 3.28, 4.68, 4.55, arrow=False)
    add_line(slide, 4.68, 4.55, 5.45, 4.55)
    add_line(slide, 6.70, 4.55, 7.05, 4.55)

    arcc = add_box(
        slide,
        8.55,
        4.08,
        2.75,
        1.45,
        "ARCC\nAnomaly Response-guided\nContext Calibration\n\nA_final = A_local + lambda * A_local * tanh(G_cal)\n\nA_local guides context calibration",
        "arcc",
        10,
        True,
    )
    arcc.line.width = Pt(1.6)
    add_line(slide, 8.25, 4.55, 8.55, 4.55)
    add_line(slide, 7.76, 3.33, 7.76, 4.08, arrow=False)
    add_line(slide, 7.76, 4.08, 8.55, 4.35)
    add_label(slide, 8.13, 3.83, 1.3, 0.22, "context", 8, italic=True)

    add_box(slide, 11.85, 4.38, 1.05, 0.58, "Pixel-level\nanomaly map", "green", 9)
    add_line(slide, 11.30, 4.78, 11.85, 4.68)

    # Loss.
    add_box(slide, 11.78, 0.78, 1.22, 1.12, "Training Loss\n\nLoc. loss\nImage BCE\nProto. reg.\nSyn. BCE/Dice", "pink", 8.5, True)
    add_line(slide, 12.70, 2.18, 12.70, 1.90)
    add_line(slide, 12.90, 4.66, 12.90, 1.90)
    add_line(slide, 5.60, 1.00, 11.78, 1.22, dashed=True)

    # Legend.
    add_box(slide, 1.05, 6.78, 10.9, 0.36, "", "white", 8, radius=True)
    legend_items = [
        ("Blue = frozen encoders", "blue"),
        ("Orange = trainable modules", "orange"),
        ("Purple = features", "purple"),
        ("Green = outputs", "green"),
        ("Pink = loss", "pink"),
    ]
    x = 1.35
    for label, style in legend_items:
        add_box(slide, x, 6.87, 0.22, 0.14, "", style, 6, radius=False)
        add_label(slide, x + 0.27, 6.82, 1.35, 0.24, label, 8)
        x += 1.9
    add_line(slide, 10.80, 6.94, 11.30, 6.94, dashed=True, width=1.0)
    add_label(slide, 11.33, 6.82, 1.0, 0.24, "training-only", 8)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    return OUT


if __name__ == "__main__":
    print(build())
